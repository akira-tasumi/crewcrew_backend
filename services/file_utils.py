"""
ファイル処理ユーティリティ

Excel, CSV, 画像, Word, PowerPoint, テキストファイルをテキストやBase64に変換するヘルパー関数
"""

import base64
import io
import json
import logging
from typing import Tuple

import pandas as pd
from PIL import Image

logger = logging.getLogger(__name__)


def process_excel_file(file_content: bytes, filename: str) -> Tuple[str, str]:
    """
    Excelファイル (.xlsx, .xls) を読み込み、Markdown形式のテキストに変換

    Args:
        file_content: ファイルのバイトデータ
        filename: ファイル名

    Returns:
        Tuple[str, str]: (変換されたテキスト, ファイルタイプ)
    """
    try:
        # ExcelファイルをDataFrameとして読み込み
        excel_file = io.BytesIO(file_content)

        # 全シートを読み込み
        all_sheets = pd.read_excel(excel_file, sheet_name=None, engine='openpyxl')

        result_texts = []
        for sheet_name, df in all_sheets.items():
            # 空のシートはスキップ
            if df.empty:
                continue

            # シート情報を追加
            result_texts.append(f"## シート: {sheet_name}")
            result_texts.append(f"行数: {len(df)}, 列数: {len(df.columns)}")
            result_texts.append("")

            # データが大きすぎる場合は先頭100行に制限
            if len(df) > 100:
                result_texts.append(f"(データが多いため、先頭100行を表示)")
                df = df.head(100)

            # Markdown形式で出力
            try:
                markdown_table = df.to_markdown(index=False)
                result_texts.append(markdown_table)
            except Exception:
                # to_markdownが使えない場合はCSV形式で出力
                result_texts.append(df.to_string(index=False))

            result_texts.append("")

        if not result_texts:
            return "Excelファイルにデータがありませんでした。", "excel"

        return "\n".join(result_texts), "excel"

    except Exception as e:
        logger.error(f"Excel processing error: {e}")
        raise ValueError(f"Excelファイルの読み込みに失敗しました: {str(e)}")


def process_csv_file(file_content: bytes, filename: str) -> Tuple[str, str]:
    """
    CSVファイルを読み込み、Markdown形式のテキストに変換

    Args:
        file_content: ファイルのバイトデータ
        filename: ファイル名

    Returns:
        Tuple[str, str]: (変換されたテキスト, ファイルタイプ)
    """
    try:
        # エンコーディングを自動判定（UTF-8, Shift-JIS, CP932を試行）
        csv_file = io.BytesIO(file_content)
        df = None

        for encoding in ['utf-8', 'shift-jis', 'cp932', 'latin-1']:
            try:
                csv_file.seek(0)
                df = pd.read_csv(csv_file, encoding=encoding)
                break
            except (UnicodeDecodeError, pd.errors.ParserError):
                continue

        if df is None:
            raise ValueError("CSVファイルのエンコーディングを判定できませんでした")

        result_texts = []
        result_texts.append(f"## CSVファイル: {filename}")
        result_texts.append(f"行数: {len(df)}, 列数: {len(df.columns)}")
        result_texts.append("")

        # データが大きすぎる場合は先頭100行に制限
        if len(df) > 100:
            result_texts.append(f"(データが多いため、先頭100行を表示)")
            df = df.head(100)

        # Markdown形式で出力
        try:
            markdown_table = df.to_markdown(index=False)
            result_texts.append(markdown_table)
        except Exception:
            result_texts.append(df.to_string(index=False))

        return "\n".join(result_texts), "csv"

    except Exception as e:
        logger.error(f"CSV processing error: {e}")
        raise ValueError(f"CSVファイルの読み込みに失敗しました: {str(e)}")


def process_image_file(file_content: bytes, filename: str) -> Tuple[str, str, str]:
    """
    画像ファイルをBase64エンコードし、メディアタイプを判定

    Args:
        file_content: ファイルのバイトデータ
        filename: ファイル名

    Returns:
        Tuple[str, str, str]: (Base64エンコードされた画像データ, メディアタイプ, ファイルタイプ)
    """
    try:
        # 画像を開いて検証
        image = Image.open(io.BytesIO(file_content))

        # 画像サイズを取得
        width, height = image.size

        # 画像が大きすぎる場合はリサイズ（Claude Vision APIの制限対策）
        max_size = 1568  # Claude推奨の最大サイズ
        if width > max_size or height > max_size:
            ratio = min(max_size / width, max_size / height)
            new_width = int(width * ratio)
            new_height = int(height * ratio)
            image = image.resize((new_width, new_height), Image.Resampling.LANCZOS)
            logger.info(f"Image resized from {width}x{height} to {new_width}x{new_height}")

        # RGBAをRGBに変換（JPEGで保存する場合）
        if image.mode == 'RGBA':
            background = Image.new('RGB', image.size, (255, 255, 255))
            background.paste(image, mask=image.split()[3])
            image = background

        # メディアタイプを判定
        format_lower = filename.lower()
        if format_lower.endswith('.png'):
            media_type = 'image/png'
            img_format = 'PNG'
        elif format_lower.endswith('.webp'):
            media_type = 'image/webp'
            img_format = 'WEBP'
        elif format_lower.endswith('.gif'):
            media_type = 'image/gif'
            img_format = 'GIF'
        else:
            media_type = 'image/jpeg'
            img_format = 'JPEG'

        # Base64エンコード
        buffer = io.BytesIO()
        image.save(buffer, format=img_format, quality=85)
        buffer.seek(0)
        base64_data = base64.standard_b64encode(buffer.read()).decode('utf-8')

        return base64_data, media_type, "image"

    except Exception as e:
        logger.error(f"Image processing error: {e}")
        raise ValueError(f"画像ファイルの読み込みに失敗しました: {str(e)}")


def process_text_file(file_content: bytes, filename: str) -> Tuple[str, str]:
    """
    テキストファイル (.txt, .md) を読み込み

    Args:
        file_content: ファイルのバイトデータ
        filename: ファイル名

    Returns:
        Tuple[str, str]: (テキスト内容, ファイルタイプ)
    """
    try:
        # エンコーディングを自動判定
        text = None
        for encoding in ['utf-8', 'shift-jis', 'cp932', 'latin-1']:
            try:
                text = file_content.decode(encoding)
                break
            except UnicodeDecodeError:
                continue

        if text is None:
            raise ValueError("テキストファイルのエンコーディングを判定できませんでした")

        # 長すぎる場合は切り詰め
        if len(text) > 50000:
            text = text[:50000] + "\n\n... (以下省略、全体で約{}文字)".format(len(text))

        file_ext = filename.lower().split('.')[-1] if '.' in filename else 'txt'
        result = f"## {file_ext.upper()}ファイル: {filename}\n\n{text}"

        return result, "text"

    except Exception as e:
        logger.error(f"Text file processing error: {e}")
        raise ValueError(f"テキストファイルの読み込みに失敗しました: {str(e)}")


def process_json_file(file_content: bytes, filename: str) -> Tuple[str, str]:
    """
    JSONファイルを読み込み、整形して表示

    Args:
        file_content: ファイルのバイトデータ
        filename: ファイル名

    Returns:
        Tuple[str, str]: (整形されたJSON, ファイルタイプ)
    """
    try:
        # エンコーディングを自動判定
        text = None
        for encoding in ['utf-8', 'shift-jis', 'cp932', 'latin-1']:
            try:
                text = file_content.decode(encoding)
                break
            except UnicodeDecodeError:
                continue

        if text is None:
            raise ValueError("JSONファイルのエンコーディングを判定できませんでした")

        # JSONをパースして整形
        data = json.loads(text)
        formatted = json.dumps(data, indent=2, ensure_ascii=False)

        # 長すぎる場合は切り詰め
        if len(formatted) > 50000:
            formatted = formatted[:50000] + "\n\n... (以下省略)"

        result = f"## JSONファイル: {filename}\n\n```json\n{formatted}\n```"

        return result, "json"

    except json.JSONDecodeError as e:
        logger.error(f"JSON parsing error: {e}")
        raise ValueError(f"JSONファイルの解析に失敗しました: {str(e)}")
    except Exception as e:
        logger.error(f"JSON file processing error: {e}")
        raise ValueError(f"JSONファイルの読み込みに失敗しました: {str(e)}")


def process_xml_file(file_content: bytes, filename: str) -> Tuple[str, str]:
    """
    XMLファイルを読み込み

    Args:
        file_content: ファイルのバイトデータ
        filename: ファイル名

    Returns:
        Tuple[str, str]: (XML内容, ファイルタイプ)
    """
    try:
        # エンコーディングを自動判定
        text = None
        for encoding in ['utf-8', 'shift-jis', 'cp932', 'latin-1']:
            try:
                text = file_content.decode(encoding)
                break
            except UnicodeDecodeError:
                continue

        if text is None:
            raise ValueError("XMLファイルのエンコーディングを判定できませんでした")

        # 長すぎる場合は切り詰め
        if len(text) > 50000:
            text = text[:50000] + "\n\n... (以下省略)"

        result = f"## XMLファイル: {filename}\n\n```xml\n{text}\n```"

        return result, "xml"

    except Exception as e:
        logger.error(f"XML file processing error: {e}")
        raise ValueError(f"XMLファイルの読み込みに失敗しました: {str(e)}")


def process_word_file(file_content: bytes, filename: str) -> Tuple[str, str]:
    """
    Word文書 (.docx) を読み込み、テキストに変換

    Args:
        file_content: ファイルのバイトデータ
        filename: ファイル名

    Returns:
        Tuple[str, str]: (抽出されたテキスト, ファイルタイプ)
    """
    try:
        from docx import Document

        doc = Document(io.BytesIO(file_content))

        result_texts = []
        result_texts.append(f"## Word文書: {filename}")
        result_texts.append("")

        # 段落を抽出
        for para in doc.paragraphs:
            if para.text.strip():
                # 見出しスタイルを検出
                if para.style and para.style.name.startswith('Heading'):
                    level = para.style.name.replace('Heading ', '')
                    try:
                        level_num = int(level)
                        result_texts.append(f"{'#' * (level_num + 1)} {para.text}")
                    except ValueError:
                        result_texts.append(f"### {para.text}")
                else:
                    result_texts.append(para.text)

        # テーブルを抽出
        for table in doc.tables:
            result_texts.append("")
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                result_texts.append("| " + " | ".join(cells) + " |")
            result_texts.append("")

        text = "\n".join(result_texts)

        # 長すぎる場合は切り詰め
        if len(text) > 50000:
            text = text[:50000] + "\n\n... (以下省略)"

        return text, "word"

    except Exception as e:
        logger.error(f"Word file processing error: {e}")
        raise ValueError(f"Word文書の読み込みに失敗しました: {str(e)}")


def process_powerpoint_file(file_content: bytes, filename: str) -> Tuple[str, str]:
    """
    PowerPointファイル (.pptx) を読み込み、テキストに変換

    Args:
        file_content: ファイルのバイトデータ
        filename: ファイル名

    Returns:
        Tuple[str, str]: (抽出されたテキスト, ファイルタイプ)
    """
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_content))

        result_texts = []
        result_texts.append(f"## PowerPointプレゼンテーション: {filename}")
        result_texts.append(f"スライド数: {len(prs.slides)}")
        result_texts.append("")

        for slide_num, slide in enumerate(prs.slides, 1):
            result_texts.append(f"### スライド {slide_num}")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    result_texts.append(shape.text)

            result_texts.append("")

        text = "\n".join(result_texts)

        # 長すぎる場合は切り詰め
        if len(text) > 50000:
            text = text[:50000] + "\n\n... (以下省略)"

        return text, "powerpoint"

    except Exception as e:
        logger.error(f"PowerPoint file processing error: {e}")
        raise ValueError(f"PowerPointファイルの読み込みに失敗しました: {str(e)}")


def get_file_type(filename: str) -> str:
    """
    ファイル名から種類を判定

    Args:
        filename: ファイル名

    Returns:
        str: ファイルタイプ
    """
    lower = filename.lower()

    if lower.endswith(('.xlsx', '.xls')):
        return 'excel'
    elif lower.endswith('.csv'):
        return 'csv'
    elif lower.endswith(('.png', '.jpg', '.jpeg', '.webp', '.gif')):
        return 'image'
    elif lower.endswith('.pdf'):
        return 'pdf'
    elif lower.endswith(('.txt', '.md')):
        return 'text'
    elif lower.endswith('.json'):
        return 'json'
    elif lower.endswith('.xml'):
        return 'xml'
    elif lower.endswith('.docx'):
        return 'word'
    elif lower.endswith('.pptx'):
        return 'powerpoint'
    else:
        return 'unknown'


def process_file(file_content: bytes, filename: str) -> dict:
    """
    ファイルを種類に応じて処理

    Args:
        file_content: ファイルのバイトデータ
        filename: ファイル名

    Returns:
        dict: 処理結果
            - type: ファイルタイプ
            - text: テキスト変換結果（Excel/CSV/テキスト/Word/PowerPoint）
            - base64: Base64データ（画像）
            - media_type: メディアタイプ（画像）
    """
    file_type = get_file_type(filename)

    if file_type == 'excel':
        text, _ = process_excel_file(file_content, filename)
        return {
            'type': 'excel',
            'filename': filename,
            'text': text,
        }

    elif file_type == 'csv':
        text, _ = process_csv_file(file_content, filename)
        return {
            'type': 'csv',
            'filename': filename,
            'text': text,
        }

    elif file_type == 'image':
        base64_data, media_type, _ = process_image_file(file_content, filename)
        return {
            'type': 'image',
            'filename': filename,
            'base64': base64_data,
            'media_type': media_type,
        }

    elif file_type == 'pdf':
        # PDFは既存の処理を使用（main.pyで処理）
        return {
            'type': 'pdf',
            'filename': filename,
            'content': file_content,
        }

    elif file_type == 'text':
        text, _ = process_text_file(file_content, filename)
        return {
            'type': 'text',
            'filename': filename,
            'text': text,
        }

    elif file_type == 'json':
        text, _ = process_json_file(file_content, filename)
        return {
            'type': 'json',
            'filename': filename,
            'text': text,
        }

    elif file_type == 'xml':
        text, _ = process_xml_file(file_content, filename)
        return {
            'type': 'xml',
            'filename': filename,
            'text': text,
        }

    elif file_type == 'word':
        text, _ = process_word_file(file_content, filename)
        return {
            'type': 'word',
            'filename': filename,
            'text': text,
        }

    elif file_type == 'powerpoint':
        text, _ = process_powerpoint_file(file_content, filename)
        return {
            'type': 'powerpoint',
            'filename': filename,
            'text': text,
        }

    else:
        raise ValueError(f"サポートされていないファイル形式です: {filename}")
